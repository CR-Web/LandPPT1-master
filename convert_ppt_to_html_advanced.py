#!/usr/bin/env python3
"""
Advanced PPT to HTML converter that creates business blue theme with proper styling
"""

import os
import sys
from pptx import Presentation
from bs4 import BeautifulSoup

# Add src to Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

def get_business_blue_styles():
    """
    Get business blue theme styles
    """
    return {
        'background_color': '#f8fafc',
        'title_box': {
            'background_color': '#1a365d',
            'font_color': '#ffffff',
            'font_size': 40,
            'padding': '30px 60px',
            'border_radius': '10px',
            'box_shadow': '0 4px 12px rgba(0,0,0,0.15)',
            'text_align': 'center'
        },
        'subtitle_box': {
            'background_color': 'rgba(26, 54, 93, 0.1)',
            'font_color': '#1a365d',
            'font_size': 28,
            'padding': '20px 40px',
            'border_radius': '8px',
            'box_shadow': '0 2px 8px rgba(0,0,0,0.1)',
            'text_align': 'center',
            'border_left': '4px solid #1a365d',
            'border_right': '4px solid #1a365d'
        },
        'content_box': {
            'background_color': 'rgba(248, 250, 252, 1)',
            'font_color': '#333333',
            'font_size': 20,
            'padding': '30px',
            'border_radius': '8px',
            'box_shadow': '0 2px 8px rgba(0,0,0,0.1)',
            'text_align': 'center',
            'border': '1px solid #e2e8f0'
        },
        'info_box': {
            'background_color': 'rgba(26, 54, 93, 0.05)',
            'font_color': '#4a5568',
            'font_size': 18,
            'padding': '20px',
            'border_radius': '8px',
            'text_align': 'center',
            'border': '1px solid rgba(26, 54, 93, 0.2)'
        }
    }

def convert_ppt_to_html(ppt_path, output_html_path):
    """
    Convert PPT file to HTML format with business blue theme
    """
    try:
        print(f"📁 Converting PPT to HTML with business blue theme: {ppt_path}")
        
        # Load presentation
        prs = Presentation(ppt_path)
        print(f"📊 Found {len(prs.slides)} slides")
        
        # Get business blue styles
        styles = get_business_blue_styles()
        
        # Prepare HTML structure
        html_parts = []
        html_parts.append('<!DOCTYPE html>\n<html lang="zh-CN">\n<head>\n    <meta charset="UTF-8">\n    <meta name="viewport" content="width=device-width, initial-scale=1.0">\n    <title>Business Blue PPT Template</title>\n    <style>')
        
        # Add base styles with business blue theme
        html_parts.append('''
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: Arial, sans-serif;
            background-color: #f0f2f5;
            padding: 20px;
        }
        
        .slide {
            width: 1280px;
            height: 720px;
            margin: 0 auto 40px;
            background-color: #f8fafc;
            border-radius: 8px;
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
            overflow: hidden;
            position: relative;
            display: flex;
            align-items: center;
            justify-content: center;
        }
        
        .slide-content {
            width: 100%;
            height: 100%;
            padding: 80px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
        }
        
        .title-box {
            background-color: #1a365d;
            color: white;
            padding: 30px 60px;
            border-radius: 10px;
            margin-bottom: 40px;
            box-shadow: 0 4px 12px rgba(0,0,0,0.15);
            text-align: center;
            max-width: 90%;
        }
        
        .title-box h1 {
            font-size: 40px;
            font-weight: bold;
            margin: 0;
        }
        
        .subtitle-box {
            background-color: rgba(26, 54, 93, 0.1);
            color: #1a365d;
            padding: 20px 40px;
            border-radius: 8px;
            margin-bottom: 30px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            text-align: center;
            max-width: 80%;
            border-left: 4px solid #1a365d;
            border-right: 4px solid #1a365d;
        }
        
        .subtitle-box h2 {
            font-size: 28px;
            font-weight: bold;
            margin: 0;
        }
        
        .content-box {
            background-color: rgba(248, 250, 252, 1);
            color: #333333;
            padding: 30px;
            border-radius: 8px;
            margin-bottom: 20px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            text-align: center;
            max-width: 80%;
            border: 1px solid #e2e8f0;
        }
        
        .content-box p {
            font-size: 20px;
            line-height: 1.5;
            margin: 0;
        }
        
        .info-box {
            background-color: rgba(26, 54, 93, 0.05);
            color: #4a5568;
            padding: 20px;
            border-radius: 8px;
            margin-bottom: 15px;
            text-align: center;
            max-width: 60%;
            border: 1px solid rgba(26, 54, 93, 0.2);
        }
        
        .info-box p {
            font-size: 18px;
            margin: 0;
        }
        
        ul {
            margin-left: 30px;
            margin-bottom: 20px;
        }
        
        li {
            font-size: 20px;
            line-height: 1.6;
            color: #555;
            margin-bottom: 10px;
        }
        
        .placeholder {
            background-color: #e0e0e0;
            border: 2px dashed #ccc;
            display: flex;
            align-items: center;
            justify-content: center;
            color: #999;
            font-size: 18px;
            height: 300px;
            width: 100%;
            max-width: 600px;
            margin: 20px 0;
        }
        ''')
        
        html_parts.append('</style>\n</head>\n<body>')
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            print(f"📄 Processing slide {i+1}")
            
            # Start slide
            html_parts.append(f'\n    <div class="slide" id="slide-{i+1}">\n        <div class="slide-content">')
            
            # Process shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    text = shape.text_frame.text.strip()
                    lines = text.split('\n')
                    
                    # Skip empty shapes
                    if not text:
                        continue
                    
                    # Determine shape type
                    is_title = False
                    is_subtitle = False
                    is_content = False
                    is_info = False
                    
                    # Check if this is a title
                    if any(keyword in text for keyword in ['标题', 'TITLE', 'title']):
                        is_title = True
                    elif any(keyword in text for keyword in ['subtitle', 'SUBTITLE']):
                        is_subtitle = True
                    elif len(lines) == 1 and len(lines[0]) < 30:
                        # Likely a title
                        is_title = True
                    elif len(text) > 100:
                        # Long content
                        is_content = True
                    else:
                        # Info text
                        is_info = True
                    
                    # Add appropriate box
                    if is_title:
                        html_parts.append('            <div class="title-box">')
                        html_parts.append(f'                <h1>{text}</h1>')
                        html_parts.append('            </div>')
                    elif is_subtitle:
                        html_parts.append('            <div class="subtitle-box">')
                        html_parts.append(f'                <h2>{text}</h2>')
                        html_parts.append('            </div>')
                    elif is_content:
                        html_parts.append('            <div class="content-box">')
                        for line in lines:
                            line = line.strip()
                            if line:
                                html_parts.append(f'                <p>{line}</p>')
                        html_parts.append('            </div>')
                    else:
                        html_parts.append('            <div class="info-box">')
                        html_parts.append(f'                <p>{text}</p>')
                        html_parts.append('            </div>')
                elif hasattr(shape, 'picture'):
                    # Image shape
                    html_parts.append('            <div class="placeholder">')
                    html_parts.append('                [Image: Please add image here]')
                    html_parts.append('            </div>')
            
            html_parts.append('        </div>\n    </div>')
        
        html_parts.append('\n</body>\n</html>')
        
        # Write HTML file
        html_content = ''.join(html_parts)
        
        # Clean up HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        clean_html = soup.prettify()
        
        with open(output_html_path, 'w', encoding='utf-8') as f:
            f.write(clean_html)
        
        print(f"✅ Advanced HTML file created: {output_html_path}")
        print(f"📋 First slide preview available at: {output_html_path}#slide-1")
        
        return True
        
    except Exception as e:
        print(f"❌ Error converting PPT to HTML: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """
    Main function
    """
    ppt_path = "e:\\workandstudy\\LandPPT-master\\src\\ppt\\business_blue_01.pptx"
    output_html_path = "e:\\workandstudy\\LandPPT-master\\src\\ppt\\business_blue_01_advanced.html"
    
    print("🚀 Starting Advanced PPT to HTML conversion...")
    print(f"📁 Input PPT: {ppt_path}")
    print(f"📄 Output HTML: {output_html_path}")
    print("=" * 60)
    
    success = convert_ppt_to_html(ppt_path, output_html_path)
    
    if success:
        print("\n🎉 Conversion completed successfully!")
        print("📌 The HTML file now includes:")
        print("   • Business blue theme with proper colors")
        print("   • Title boxes with background colors")
        print("   • Content boxes with proper styling")
        print("   • Info boxes for additional content")
        print("💡 When importing to frontend, the template will preview correctly")
    else:
        print("\n❌ Conversion failed")

if __name__ == "__main__":
    main()