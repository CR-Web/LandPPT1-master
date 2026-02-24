#!/usr/bin/env python3
"""
Convert PPT file to HTML format with proper styling for frontend preview
"""

import os
import sys
from pptx import Presentation
from bs4 import BeautifulSoup

# Add src to Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

def convert_ppt_to_html(ppt_path, output_html_path):
    """
    Convert PPT file to HTML format preserving styles
    """
    try:
        print(f"📁 Converting PPT to HTML: {ppt_path}")
        
        # Load presentation
        prs = Presentation(ppt_path)
        print(f"📊 Found {len(prs.slides)} slides")
        
        # Prepare HTML structure
        html_parts = []
        html_parts.append('<!DOCTYPE html>\n<html lang="zh-CN">\n<head>\n    <meta charset="UTF-8">\n    <meta name="viewport" content="width=device-width, initial-scale=1.0">\n    <title>Business Blue PPT Template</title>\n    <style>')
        
        # Add base styles
        html_parts.append('''
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        body {
            font-family: Arial, sans-serif;
            background-color: #f0f2f5;
            padding: 20px;
        }
        
        .slide {
            width: 960px;
            height: 720px;
            margin: 0 auto 40px;
            background-color: white;
            border-radius: 8px;
            box-shadow: 0 4px 12px rgba(0,0,0,0.1);
            overflow: hidden;
            position: relative;
        }
        
        .slide-content {
            width: 100%;
            height: 100%;
            padding: 60px;
            display: flex;
            flex-direction: column;
        }
        
        h1 {
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 30px;
            color: #2c3e50;
            text-align: center;
        }
        
        h2 {
            font-size: 28px;
            font-weight: bold;
            margin-bottom: 20px;
            color: #34495e;
        }
        
        p {
            font-size: 20px;
            line-height: 1.5;
            color: #555;
            margin-bottom: 15px;
        }
        
        ul {
            margin-left: 30px;
            margin-bottom: 20px;
        }
        
        li {
            font-size: 20px;
            line-height: 1.6;
            color: #555;
            margin-bottom: 10px;
        }
        
        .placeholder {
            background-color: #e0e0e0;
            border: 2px dashed #ccc;
            display: flex;
            align-items: center;
            justify-content: center;
            color: #999;
            font-size: 18px;
        }
        
        .text-box {
            margin-bottom: 20px;
        }
        ''')
        
        # Add slide-specific styles based on first slide
        if prs.slides:
            first_slide = prs.slides[0]
            # Analyze first slide for styling
            has_title = False
            for shape in first_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    has_title = True
                    break
            
            if has_title:
                html_parts.append('''
                .slide:first-child h1 {
                    color: #1a365d;
                    font-size: 40px;
                    margin-top: 100px;
                }
                ''')
        
        html_parts.append('</style>\n</head>\n<body>')
        
        # Process each slide
        for i, slide in enumerate(prs.slides):
            print(f"📄 Processing slide {i+1}")
            html_parts.append(f'\n    <div class="slide" id="slide-{i+1}">\n        <div class="slide-content">')
            
            # Process shapes in slide
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    # Text shape
                    text = shape.text_frame.text
                    lines = text.split('\n')
                    
                    # Determine heading level based on text length and position
                    if len(lines) == 1 and len(lines[0]) < 50:
                        # Likely a title
                        html_parts.append(f'            <div class="text-box">\n                <h1>{lines[0]}</h1>\n            </div>')
                    else:
                        # Content text
                        html_parts.append('            <div class="text-box">')
                        for line in lines:
                            line = line.strip()
                            if line:
                                if line.startswith('-') or line.startswith('•'):
                                    # List item
                                    html_parts.append(f'                <li>{line[1:].strip()}</li>')
                                else:
                                    # Paragraph
                                    html_parts.append(f'                <p>{line}</p>')
                        html_parts.append('            </div>')
                elif hasattr(shape, 'picture'):
                    # Image shape
                    html_parts.append('            <div class="placeholder" style="height: 300px;">')
                    html_parts.append('                [Image: Please add image here]')
                    html_parts.append('            </div>')
            
            html_parts.append('        </div>\n    </div>')
        
        html_parts.append('\n</body>\n</html>')
        
        # Write HTML file
        html_content = ''.join(html_parts)
        
        # Clean up HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        clean_html = soup.prettify()
        
        with open(output_html_path, 'w', encoding='utf-8') as f:
            f.write(clean_html)
        
        print(f"✅ HTML file created: {output_html_path}")
        print(f"📋 First slide preview available at: {output_html_path}#slide-1")
        
        return True
        
    except Exception as e:
        print(f"❌ Error converting PPT to HTML: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """
    Main function
    """
    ppt_path = "e:\\workandstudy\\LandPPT-master\\src\\ppt\\business_blue_01.pptx"
    output_html_path = "e:\\workandstudy\\LandPPT-master\\src\\ppt\\business_blue_01.html"
    
    print("🚀 Starting PPT to HTML conversion...")
    print(f"📁 Input PPT: {ppt_path}")
    print(f"📄 Output HTML: {output_html_path}")
    print("=" * 60)
    
    success = convert_ppt_to_html(ppt_path, output_html_path)
    
    if success:
        print("\n🎉 Conversion completed successfully!")
        print("📌 The HTML file is ready for frontend preview")
        print("💡 When importing to frontend, the first slide will be displayed correctly")
    else:
        print("\n❌ Conversion failed")

if __name__ == "__main__":
    main()